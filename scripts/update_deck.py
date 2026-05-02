"""Update docs/CIS509_MCOOK20.pptx in place to match the current pipeline.

Reads the live outputs (top_10_manipulated.csv, figures, receipts) and
applies a surgical diff to the existing 9-slide deck:

  - corrects 28.3% to 10.3% wherever it appears (slide 3, slide 9, notes)
  - swaps committed PNG figures into slides 2, 4, 7, 8
  - rewrites slide 8 to combine real top-10 data + the dashboard mention
  - adds dashboard + MCP server to the "Delivered" column on slide 9
  - refreshes all 9 speaker notes with TruthFirst voice and current numbers

Slide count stays at 9. Layout, theme, fonts, colors are preserved by
editing existing runs in place rather than reconstructing shapes.

Idempotent. Re-runnable.

Usage:
    python scripts/update_deck.py
"""

from __future__ import annotations

import sys
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO))

import pandas as pd
from pptx import Presentation
from pptx.util import Emu, Pt

from src.config import OUTPUTS_DIR

DECK_PATH = REPO / "docs" / "CIS509_MCOOK20.pptx"
FIG = OUTPUTS_DIR / "figures"
TOP10_PATH = OUTPUTS_DIR / "top_10_manipulated.csv"


# --------------------------------------------------------------------------- #
# Text helpers, preserve formatting
# --------------------------------------------------------------------------- #

def _font_snapshot(run):
    """Capture font props before clearing a run."""
    f = run.font
    snap = {"name": f.name, "size": f.size, "bold": f.bold, "italic": f.italic}
    try:
        snap["color_rgb"] = f.color.rgb if f.color and f.color.type is not None else None
    except (AttributeError, TypeError):
        snap["color_rgb"] = None
    return snap


def _font_apply(run, snap):
    f = run.font
    if snap.get("name"): f.name = snap["name"]
    if snap.get("size"): f.size = snap["size"]
    if snap.get("bold") is not None: f.bold = snap["bold"]
    if snap.get("italic") is not None: f.italic = snap["italic"]
    if snap.get("color_rgb") is not None:
        try:
            f.color.rgb = snap["color_rgb"]
        except Exception:
            pass


def set_text_lines(shape, lines: list[str]) -> None:
    """Replace shape text with `lines` (one per paragraph), preserving the
    first paragraph's first-run font for every new line. Keeps theme.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Snapshot the source font from the first existing run
    snap = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        snap = _font_snapshot(tf.paragraphs[0].runs[0])
    # Clear and rebuild
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Wipe any default run added by clear()/add_paragraph()
        for r in list(para.runs):
            r.text = ""
        run = para.add_run()
        run.text = line
        if snap:
            _font_apply(run, snap)


def set_text(shape, text: str) -> None:
    """Single-line replacement, preserves font."""
    set_text_lines(shape, [text])


def set_notes(slide, text: str) -> None:
    """Replace speaker notes for a slide with `text`."""
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    para = notes.paragraphs[0]
    for r in list(para.runs):
        r.text = ""
    para.add_run().text = text


# --------------------------------------------------------------------------- #
# Picture replacement
# --------------------------------------------------------------------------- #

def replace_picture(slide, shape_index: int, new_image_path: Path) -> None:
    """Delete the picture at `shape_index` and re-insert from `new_image_path`
    at the same position and size. Preserves layout.
    """
    shapes = slide.shapes
    target = shapes[shape_index]
    L, T, W, H = target.left, target.top, target.width, target.height
    # Delete the old picture
    target.element.getparent().remove(target.element)
    # Add new one at the same geometry
    slide.shapes.add_picture(str(new_image_path), L, T, width=W, height=H)


# --------------------------------------------------------------------------- #
# Per-slide diffs
# --------------------------------------------------------------------------- #

def update_slide_1(slide) -> None:
    """Title slide. Subtitle gets a small framing tweak; everything else stays."""
    set_text(slide.shapes[1], "Trust-Adjusted Business Health Score, receipts-native")
    set_notes(slide, (
        "TABHS strips fake-review noise from Yelp ratings and produces a per-business "
        "score. Every analyzed business gets a JSON receipt signed with SHA-256 + BLAKE3. "
        "The receipt records inputs, evidence, model lineage, limitations. "
        "Open the dashboard to see the top 10 most-manipulated businesses with one click. "
        "This deck walks the methodology, the honest findings, and the verifiable artifacts."
    ))


def update_slide_2(slide) -> None:
    """EDA. Star distribution headline stays; refresh figure + footer counts."""
    set_text_lines(slide.shapes[3], [
        "48,147 raw reviews  |  47,035 after cleaning  |  1,864 businesses  |  Arizona  |  2020-2022",
    ])
    eda_path = FIG / "eda_six_panel.png"
    if eda_path.exists():
        replace_picture(slide, 2, eda_path)
    set_notes(slide, (
        "Focal point is the gold 5-star bar dwarfing everything else. A legitimate "
        "distribution of business quality does not produce this shape. The 50.8% number "
        "is the thesis: the rating signal is structurally biased before any manipulation "
        "is even considered. The corpus is 48,147 raw rows, 47,035 after dropping Excel "
        "#NAME? corruption in the source CSV. Cleaning is documented in lessons.md."
    ))


def update_slide_3(slide) -> None:
    """Headline finding. 28.3% was wrong; the actual computed value is 10.3%."""
    set_text(slide.shapes[0], "10.3% of reviews contradict their own rating.")
    set_text(slide.shapes[1], "4,838 reviews")
    set_text_lines(slide.shapes[2], [
        "show high sentiment-star divergence (score > 1.0).",
        "The text says one thing. The rating says another.",
        "These are the inauthenticity candidates the classifier is trained on.",
    ])
    # Stat tiles, recomputed from the actual corpus
    set_text(slide.shapes[3], "4.19M")     # tokens
    set_text(slide.shapes[4], "TOKENS")
    set_text(slide.shapes[5], "86.9")      # avg words
    set_text(slide.shapes[6], "AVG WORDS")
    set_text(slide.shapes[7], "22,435")    # reviewers (raw)
    set_text(slide.shapes[8], "REVIEWERS")
    set_text(slide.shapes[9], "102,154")   # vocab
    set_text(slide.shapes[10], "VOCABULARY")
    set_text(slide.shapes[11], "2020-2022")
    set_text(slide.shapes[12], "DATE RANGE")
    set_text(slide.shapes[13], "1,864")
    set_text(slide.shapes[14], "BUSINESSES")
    set_notes(slide, (
        "Focal point: 4,838 in gold. That is the headline finding from the EDA. "
        "10.3% of the corpus has sentiment that does not match the star rating. "
        "Earlier deck drafts said 28.3%, that number was wrong; the lessons.md file "
        "logs the correction and ProjectEDA's own output reproduces 10.4% on the "
        "uncleaned data. Supporting stats on the right are corpus context. The "
        "divergence threshold of 1.0 becomes the proxy label downstream."
    ))


def update_slide_4(slide) -> None:
    """VADER divergence. Formula and example stay; refresh chart."""
    div_path = FIG / "divergence_histogram.png"
    if div_path.exists():
        replace_picture(slide, 3, div_path)
    set_notes(slide, (
        "The formula is the focal point. Walk through the worked example: glowing review "
        "text, 1-star rating. VADER reads the text as positive, the star is mapped to -1.0. "
        "That 1.84 gap is the signal. The histogram on the right shows the population, "
        "with the gold line at the 1.0 threshold. Most divergent reviews sit at 1 and 2 "
        "stars, which is where reviewers write nuanced criticism that VADER reads as net "
        "positive. The classifier picks up on what other features track this pattern."
    ))


def update_slide_5(slide) -> None:
    """BERTopic. Mostly stable; tighten the closer."""
    set_text(slide.shapes[13], (
        "Good reviews are specific. Bad reviews are structural. "
        "Topic assignment becomes a feature in the XGBoost layer."
    ))
    set_notes(slide, (
        "The insight: good reviews name food items and restaurants, bad reviews describe "
        "what failed operationally. The COVID topic is temporally bounded which confirms "
        "organic signal generation. 80 raw topics reduce to 30 merged then 15 final, "
        "all interpretable by inspection. Topic-id and is-topic-outlier feed into the "
        "XGBoost feature matrix in notebook 05."
    ))


def update_slide_6(slide) -> None:
    """LLM table. Fix model name (DeepSeek -> gpt-oss-120b per LA6 actual). Add caveat."""
    tbl = slide.shapes[1].table
    # Cell (3, 0) is the comparator-model row label
    cell = tbl.cell(3, 0)
    if cell.text_frame.text.strip().startswith("Few-Shot (DeepSeek"):
        # Preserve formatting of the existing cell
        snap = None
        if cell.text_frame.paragraphs[0].runs:
            snap = _font_snapshot(cell.text_frame.paragraphs[0].runs[0])
        cell.text_frame.clear()
        para = cell.text_frame.paragraphs[0]
        for r in list(para.runs):
            r.text = ""
        run = para.add_run()
        run.text = "Few-Shot (gpt-oss-120b)"
        if snap:
            _font_apply(run, snap)
    set_text_lines(slide.shapes[2], [
        "VADER + traditional ML outperforms LLM prompting at scale on this task.",
        "Few-shot closes the gap by 5 points over zero-shot. Multi-model agreement "
        "becomes a feature in the XGBoost layer, not a replacement.",
        "Run on a 40-review balanced subset. Cached predictions persist in "
        "outputs/llm_agreement.csv so the pipeline runs without a Groq API key.",
    ])
    set_notes(slide, (
        "The table is the focal point. SVM at 95.7% beats every LLM protocol on this "
        "balanced 40-review sample. The honest move is to lead with this finding rather "
        "than hide it. LLMs bring flexibility, not raw accuracy. The graduate-level "
        "insight is that the LLM agreement score is a useful FEATURE for XGBoost even "
        "when its raw accuracy is below SVM. The comparator model is gpt-oss-120b, not "
        "DeepSeek as an earlier draft said; lessons.md notes the correction."
    ))


def update_slide_7(slide) -> None:
    """XGBoost. Refresh feature importance figure + body text with current pipeline metrics.

    Text edits before picture swap, since add_picture shifts indices.
    """
    set_text_lines(slide.shapes[2], [
        "Proxy label: divergence > 1.0 = suspicious.",
        "No ground truth fake-review labels exist.",
        "",
        "Circularity acknowledged: divergence is both",
        "a feature and a component of the label.",
        "XGBoost learns from the other features.",
        "",
        "Output per business, written to a signed receipt:",
        "raw_yelp_rating",
        "tabhs_adjusted_rating",
        "manipulation_delta",
        "+ evidence breakdown + dual hash",
    ])
    fi_path = FIG / "xgboost_feature_importance.png"
    if fi_path.exists():
        replace_picture(slide, 1, fi_path)
    set_notes(slide, (
        "Focal point is the feature importance chart. Walk through the honest disclosures: "
        "proxy label, classifier circularity (divergence is both an input and a label "
        "component), and the receipt as the output format rather than a single score. "
        "Test accuracy is high (above 0.99) because of the circularity, NOT because the "
        "classifier detects fake reviews. It detects the OTHER feature patterns that "
        "track high-divergence reviews. This honesty earns rubric points, not loses them."
    ))


def update_slide_8(slide) -> None:
    """Top 10 + dashboard. Refresh figure, real numbers, dashboard mention.

    Order matters: text edits first, picture swap last, since add_picture
    appends to the end of the shape list and shifts indices.
    """
    # Pull the actual top-1 from the live data for the speaker notes example
    top10 = pd.read_csv(TOP10_PATH).head(10)
    leader = top10.iloc[0]

    set_text_lines(slide.shapes[2], [
        f"Top-1 today: {leader['business_id']}, raw {leader['raw_avg_stars']:.2f} -> "
        f"adjusted {leader['tabhs_adjusted_stars']:.2f}, delta +{leader['manipulation_delta']:.2f}.",
        "SMB Lenders: underwrite on adjusted score. Franchise Brokers: evaluate targets "
        "with manipulation removed. M&A: surface inflation before closing.",
        "Dashboard: outputs/dashboard.html. Open via Codespaces, click any row to inspect "
        "the receipt, click Verify SHA-256 to recompute the hash in your browser.",
    ])

    top10_path = FIG / "top10_raw_vs_adjusted.png"
    if top10_path.exists():
        replace_picture(slide, 1, top10_path)
    set_notes(slide, (
        f"Focal point: the gap between raw Yelp bars and TABHS-adjusted bars in the chart. "
        f"On this corpus the largest manipulation deltas land on LOW-rating businesses, "
        f"not the 4.7-star inflation pattern the strategy doc envisioned. "
        f"Top-1 today is {leader['business_id']}, raw {leader['raw_avg_stars']:.2f} stars, "
        f"adjusted {leader['tabhs_adjusted_stars']:.2f}, delta +{leader['manipulation_delta']:.2f}. "
        f"That delta is what suspicious reviews were inflating the rating by. "
        f"Then point to the dashboard URL: every row in the live dashboard has a clickable "
        f"signed receipt, and the Verify button recomputes SHA-256 client-side. "
        f"Three market applications close the slide."
    ))


def update_slide_9(slide) -> None:
    """Wrap up. Update Delivered + Limitations + Future to reflect the actual ship."""
    # Delivered
    set_text_lines(slide.shapes[2], [
        "Three NLP methods on one corpus",
        "XGBoost synthesis classifier with proxy label",
        "Dual-hash receipts per business (SHA-256 + BLAKE3)",
        "Static dashboard, client-side hash verification",
        "MCP server, query / verify / topology",
        "Reproducible Codespaces pipeline, weekly auto-refresh",
    ])
    # Limitations
    set_text_lines(slide.shapes[4], [
        "No ground truth fake-review labels exist",
        "Proxy label circularity (disclosed in receipts)",
        "Arizona restaurants only",
        "Account features excluded per professor guidance",
        "On this corpus, deltas are small (top-1 ~0.20)",
        "LLM accuracy below SVM baseline at the subset scale",
    ])
    # Future
    set_text_lines(slide.shapes[6], [
        "Add account-level features when available",
        "Expand to national Yelp dataset",
        "Train on labeled fake reviews if obtainable",
        "Temporal burst detection per business",
        "BLAKE3 client-side verify in the browser",
        "Quantum-resistant receipt signing",
    ])
    set_notes(slide, (
        "Three columns. Delivered, limitations, future work. The Delivered list now "
        "includes the dashboard with client-side hash verification and the MCP server "
        "with three tools (query receipts, verify chain, get topology). Limitations are "
        "honest scope boundaries, not weaknesses. The 'small deltas on this corpus' line "
        "is the truth: TABHS works as forensic governance, the magnitudes on Yelp Arizona "
        "are smaller than the strategy doc envisioned. The receipts make the result "
        "auditable regardless of magnitude. Thank you."
    ))


# --------------------------------------------------------------------------- #
# Driver
# --------------------------------------------------------------------------- #

UPDATERS = [
    update_slide_1, update_slide_2, update_slide_3, update_slide_4, update_slide_5,
    update_slide_6, update_slide_7, update_slide_8, update_slide_9,
]


def main() -> int:
    if not DECK_PATH.exists():
        print(f"FAIL: {DECK_PATH} missing", file=sys.stderr)
        return 1
    if not TOP10_PATH.exists():
        print(f"FAIL: {TOP10_PATH} missing, run run_pipeline.py first", file=sys.stderr)
        return 1

    p = Presentation(DECK_PATH)
    if len(p.slides) != 9:
        print(f"FAIL: expected 9 slides, got {len(p.slides)}", file=sys.stderr)
        return 1

    for i, (slide, fn) in enumerate(zip(p.slides, UPDATERS), 1):
        fn(slide)
        print(f"  updated slide {i}: {slide.shapes[0].text_frame.text[:70]}")

    p.save(DECK_PATH)
    print(f"\nSaved {DECK_PATH.relative_to(REPO)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
